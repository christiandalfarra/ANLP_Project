#!/usr/bin/env python3
"""Alternate-style project deck (report/slides/ANLP_slides_alt.pptx).

Same content and numbers as make_slides.py, restyled as a dark editorial deck:
charcoal background, serif (Georgia) headlines, teal accent, left accent rules
instead of full title bars, and figures placed on white cards.
"""
import os
import struct
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ---- palette (dark editorial) ----
BG     = RGBColor(0x0F, 0x14, 0x1B)   # near-black slate
CARD   = RGBColor(0x1A, 0x22, 0x2E)   # panel
CARD2  = RGBColor(0x24, 0x2F, 0x3E)   # lighter panel / header cells
INK    = RGBColor(0xF4, 0xF6, 0xFA)   # primary text
MUTED  = RGBColor(0x93, 0x9F, 0xB0)   # secondary text
ACCENT = RGBColor(0x3D, 0xD6, 0xB8)   # teal
CORAL  = RGBColor(0xF2, 0x7A, 0x6E)   # problem / warning
GOLD   = RGBColor(0xF3, 0xC0, 0x59)   # secondary highlight
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)

HEAD = "Georgia"   # serif headlines
BODY = "Calibri"   # sans body

SW, SH = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width, prs.slide_height = SW, SH
BLANK = prs.slide_layouts[6]

RECT, ROUND = 1, 5  # autoshape type ids


def png_size(path):
    d = open(path, "rb").read(26)
    w, h = struct.unpack(">II", d[16:24])
    return w, h


def add_slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(RECT, 0, 0, SW, SH)   # full-bleed background
    bg.fill.solid(); bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    return s


def textbox(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    return tf


def set_run(r, text, size, color=INK, bold=False, italic=False, font=BODY):
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = font


def header(slide, n, kicker, title, sub=None):
    """Left accent rule + kicker + serif title (no full-width bar)."""
    rule = slide.shapes.add_shape(RECT, Inches(0.6), Inches(0.55), Inches(0.09), Inches(0.95))
    rule.fill.solid(); rule.fill.fore_color.rgb = ACCENT
    rule.line.fill.background()
    tf = textbox(slide, Inches(0.85), Inches(0.42), Inches(11.6), Inches(1.25))
    p = tf.paragraphs[0]
    set_run(p.add_run(), kicker.upper(), 12, ACCENT, bold=True)
    p.add_run().text = "  "
    set_run(p.add_run(), f"· {n:02d}", 12, MUTED, bold=True)
    p2 = tf.add_paragraph()
    set_run(p2.add_run(), title, 26, INK, bold=True, font=HEAD)
    if sub:
        p3 = tf.add_paragraph()
        set_run(p3.add_run(), sub, 13, MUTED)


def bullets(slide, items, x=Inches(0.85), y=Inches(1.95), w=Inches(11.7),
            h=Inches(5.0), size=17):
    tf = textbox(slide, x, y, w, h)
    first = True
    for item in items:
        if isinstance(item, str):
            item = [(item, {})]
        level = item[0][1].pop("level", 0) if item else 0
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = level
        p.space_after = Pt(9)
        # accent tick for top-level non-empty lines
        if level == 0 and item and item[0][0]:
            set_run(p.add_run(), "▸ ", size - 2, ACCENT, bold=True)
        for text, kw in item:
            sz = kw.pop("size", size - 2 * level)
            kw.setdefault("color", INK)
            set_run(p.add_run(), text, sz, **kw)
    return tf


def footer(slide, n):
    tf = textbox(slide, Inches(0.6), SH - Inches(0.5), Inches(6.0), Inches(0.32))
    p = tf.paragraphs[0]
    set_run(p.add_run(), "Summarising Football Commentaries", 10, MUTED)
    tf = textbox(slide, SW - Inches(1.1), SH - Inches(0.5), Inches(0.8), Inches(0.32))
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
    set_run(p.add_run(), f"{n:02d}", 11, ACCENT, bold=True)


def card(slide, x, y, w, h, fill=CARD, border=CARD2):
    box = slide.shapes.add_shape(ROUND, x, y, w, h)
    box.fill.solid(); box.fill.fore_color.rgb = fill
    box.line.color.rgb = border; box.line.width = Pt(1)
    return box


def figure_card(slide, path, x, y, width, pad=Inches(0.16)):
    """White rounded card sized to the image, image centred on top."""
    w, h = png_size(path)
    img_h = int(width * h / w)
    cw, ch = width + 2 * pad, img_h + 2 * pad
    card(slide, x, y, cw, ch, fill=WHITE, border=CARD2)
    slide.shapes.add_picture(path, x + pad, y + pad, width=width)
    return cw, ch


def quote_card(slide, x, y, w, h, label, text, accent=ACCENT):
    box = card(slide, x, y, w, h, fill=CARD, border=CARD2)
    # accent spine
    spine = slide.shapes.add_shape(RECT, x, y + Inches(0.12), Inches(0.06), h - Inches(0.24))
    spine.fill.solid(); spine.fill.fore_color.rgb = accent
    spine.line.fill.background()
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.24); tf.margin_right = Inches(0.18); tf.margin_top = Inches(0.12)
    p = tf.paragraphs[0]
    set_run(p.add_run(), label, 12, accent, bold=True)
    p2 = tf.add_paragraph()
    set_run(p2.add_run(), text, 13, INK, italic=True)


# ---------------------------------------------------------------- 1 · Title
s = add_slide()
# big accent block on the left
blk = s.shapes.add_shape(RECT, 0, 0, Inches(0.22), SH)
blk.fill.solid(); blk.fill.fore_color.rgb = ACCENT; blk.line.fill.background()
tf = textbox(s, Inches(0.95), Inches(1.55), Inches(11.6), Inches(0.5))
set_run(tf.paragraphs[0].add_run(),
        "APPLIED NATURAL LANGUAGE PROCESSING  ·  UNIVERSITY OF TRENTO  ·  MAY 2026",
        13, ACCENT, bold=True)
tf = textbox(s, Inches(0.9), Inches(2.2), Inches(11.7), Inches(2.6))
set_run(tf.paragraphs[0].add_run(),
        "Summarising Football Match Commentaries", 40, INK, bold=True, font=HEAD)
p = tf.add_paragraph()
set_run(p.add_run(),
        "A comparative study of prompting and fine-tuning on long-form sports audio transcripts",
        20, MUTED)
# rule
r = s.shapes.add_shape(RECT, Inches(0.95), Inches(5.15), Inches(4.2), Inches(0.035))
r.fill.solid(); r.fill.fore_color.rgb = CARD2; r.line.fill.background()
tf = textbox(s, Inches(0.9), Inches(5.35), Inches(11.7), Inches(1.0))
set_run(tf.paragraphs[0].add_run(), "Marco Panciera   ·   Christian Dalfarra", 18, INK, bold=True)

# ------------------------------------------------- 2 · Problem & questions
s = add_slide()
header(s, 2, "The task", "2.5 hours of radio → a 200-word match report")
bullets(s, [
    [("Input:  ", {"bold": True, "color": ACCENT}),
     ("Whisper transcript of live BBC Radio 5 commentary — mean ~25,000 words, "
      "noisy, fragmented, full of irrelevant detail.", {})],
    [("Output:  ", {"bold": True, "color": ACCENT}),
     ("a structured match report (~200 words): score line, key events, short tactical note.", {})],
    [("Why it is hard:  ", {"bold": True, "color": ACCENT}),
     ("inputs exceed every standard transformer context window (1,024 tokens ≈ 4% of a match); "
      "references for training are scarce.", {})],
    "",
    [("Two empirically answerable questions", {"bold": True, "color": GOLD})],
    [("1.  Chunk-and-aggregate, or a long-context encoder (LED-16384)?", {"level": 1})],
    [("2.  For each input strategy — prompting a general model, or fine-tuning on small in-domain data?",
      {"level": 1})],
], size=18)
footer(s, 2)

# ------------------------------------------------------------- 3 · Dataset
s = add_slide()
header(s, 3, "Dataset", "99 matches, 43 years, 12 competition categories",
       "BBC radio audio  →  Whisper transcription  →  GPT-4 reference summaries")
bullets(s, [
    [("Audio:  ", {"bold": True, "color": ACCENT}),
     ("BBC Radio 5 Live archive (YouTube, via yt-dlp) — ~19 GB MP3, matches from 1983 to 2026.", {})],
    [("Transcription:  ", {"bold": True, "color": ACCENT}),
     ("Whisper-small; full transcript + time-aligned segments with pitch/energy prosodic features.", {})],
    [("References:  ", {"bold": True, "color": ACCENT}),
     ("GPT-4 (ChatGPT web, search enabled) prompted per match identifier; facts grounded in "
      "web-retrieved match reports, BBC sourcing verified for every match; fixed template, ~190 words. "
      "A pragmatic, disclosed compromise — no human summaries exist for radio commentary.", {})],
    [("Split:  ", {"bold": True, "color": ACCENT}),
     ("80 train / 10 val / 9 test, stratified, committed to the repo for reproducibility.", {})],
    "",
    [("Known noise floor:  ", {"bold": True, "color": CORAL}),
     ("Whisper systematically corrupts player names (“Bergkamp” → “Burgkamp”) — "
      "an irreducible ceiling on surface-overlap metrics.", {})],
], y=Inches(2.15), size=17)
footer(s, 3)

# -------------------------------------------------- 4 · Experimental design
s = add_slide()
header(s, 4, "Design", "2 input pipelines × 2 learning paradigms",
       "3 architectures · 3 prompt styles · 11 evaluated conditions")
rows = [
    ["", "Chunk + aggregate", "Long-context (16,384 tokens)"],
    ["Prompting\n(zero / few / CoT)", "FLAN-T5-Large (512)\nBART-Large-CNN (1,024)", "LED-base-16384"],
    ["Fine-tuning", "BART as the merge step\n(“chunk-merger” recipe)", "LED on 8,192-token\nrandom windows"],
]
tbl = s.shapes.add_table(3, 3, Inches(0.85), Inches(2.15), Inches(11.6), Inches(2.9)).table
tbl.columns[0].width = Inches(2.9)
tbl.columns[1].width = Inches(4.35)
tbl.columns[2].width = Inches(4.35)
for i, row in enumerate(rows):
    for j, cell_text in enumerate(row):
        cell = tbl.cell(i, j)
        cell.fill.solid()
        hdr = i == 0 or j == 0
        cell.fill.fore_color.rgb = CARD2 if hdr else CARD
        tfc = cell.text_frame; tfc.word_wrap = True
        for k, line in enumerate(cell_text.split("\n")):
            p = tfc.paragraphs[0] if k == 0 else tfc.add_paragraph()
            set_run(p.add_run(), line, 15 if hdr else 14,
                    ACCENT if hdr else INK, bold=hdr)
bullets(s, [
    [("Chunk + aggregate:  ", {"bold": True, "color": ACCENT}),
     ("segment-aligned chunks → per-chunk summary → merge + dedup.    ", {}),
     ("Long-context:  ", {"bold": True, "color": ACCENT}),
     ("whole transcript in one LED pass.", {})],
    [("All experiments on free-tier Kaggle (1× T4); seeds fixed; every prediction saved to disk, "
      "so all metrics are recomputable post-hoc.", {})],
], y=Inches(5.35), size=15)
footer(s, 4)

# ------------------------------------------------------------- 5 · Results
s = add_slide()
header(s, 5, "Results", "Fine-tuning wins, and the gap is significant",
       "Test ROUGE-L on the clean 9-match test set · bootstrap 95% CIs (10,000 resamples)")
figure_card(s, "figures/leaderboard.png", Inches(2.7), Inches(1.95), Inches(7.6))
tf = textbox(s, Inches(0.85), Inches(6.55), Inches(11.7), Inches(0.8))
set_run(tf.paragraphs[0].add_run(),
        "Every prompting condition is significantly below both fine-tuned models (paired bootstrap; 95% CI "
        "of the difference excludes 0). The LED−BART gap (−0.0004 [−0.042, +0.035]) is NOT significant.",
        11.5, MUTED, italic=True)
footer(s, 5)

# ----------------------------------------- 6 · Finding 1: chunk-merger
s = add_slide()
header(s, 6, "Finding 1", "Match the training distribution, not the context window",
       "The “chunk-merger” fine-tuning recipe")
bullets(s, [
    [("Naive fine-tuning ", {"bold": True}),
     ("on (transcript[:1024] → reference) shows BART only the first ~5 minutes of a 2.5-hour match.", {})],
    [("Chunk-merger recipe:  ", {"bold": True, "color": ACCENT}),
     ("run pre-trained BART on every chunk → concatenate the ~30 mini-summaries (~960 tokens) → "
      "fine-tune BART to map that intermediate to the reference report. "
      "Training input now matches what the model sees at inference.", {})],
    [("Two-phase schedule: ", {"bold": True}),
     ("encoder frozen 3 epochs, then unfrozen; early stopping on val ROUGE-L.", {})],
    "",
    [("Same model, same pipeline, same test set", {"bold": True, "color": GOLD})],
    [("pre-trained BART 0.147   →   fine-tuned BART 0.248   ( +69% relative )",
      {"size": 23, "bold": True, "color": ACCENT, "level": 1})],
    [("+55% over the strongest prompting condition of any architecture (LED zero-shot, 0.160).",
      {"level": 1})],
], size=18)
footer(s, 6)

# ---------------------------------- 7 · Finding 2: ranking inversion
s = add_slide()
header(s, 7, "Finding 2", "Long context wins at prompting, ties at fine-tuning")
bullets(s, [
    [("Prompting:  ", {"bold": True, "color": ACCENT}),
     ("LED long-context (avg 0.158) beats BART chunk (0.124) and FLAN chunk (0.058) — "
      "the 16k window preserves cross-segment information that chunking discards.", {})],
    [("Fine-tuning:  ", {"bold": True, "color": ACCENT}),
     ("LED long-context is nominally ahead (0.248 vs 0.248) but the paired difference is not "
      "significant at n = 9 — the two fine-tuned pipelines cannot be separated. LED is more "
      "consistent (CI [0.230, 0.266]); BART has more variance (CI [0.208, 0.283]).", {})],
    "",
    [("Why the two pipelines are statistically tied despite architectural differences:",
      {"bold": True, "color": GOLD})],
    [("BART learns from a ~960-token pre-distilled intermediate whose features are already informative.",
      {"level": 1})],
    [("LED must learn what matters inside raw 8,192-token windows from scratch — "
      "too few pairs (N = 80) to learn long-range attention patterns reliably.", {"level": 1})],
    "",
    [("Secondary finding: ", {"bold": True}),
     ("zero-shot beats few-shot and CoT in both chunked families — examples eat the input "
      "budget; LED few-shot fell back to zero-shot (implementation gap, disclosed in the report), "
      "and LED-base is not instruction-tuned, so prompt style barely matters.", {})],
], size=17)
footer(s, 7)

# --------------------------------------- 8 · Finding 3: hallucination
s = add_slide()
header(s, 8, "Finding 3", "Models learn the format, not the facts",
       "Real test-set outputs; each quote labelled with the true result")
quote_card(s, Inches(0.85), Inches(2.15), Inches(5.75), Inches(2.25),
           "Fine-tuned BART (best model) — 2019 UCL Final (true: Liverpool 2–0 Tottenham, Madrid, no shootout)",
           "“Liverpool 2–2 Tottenham Hotspur (Liverpool win 4–3 on penalties) … 2023 UEFA "
           "Champions League Final, Wembley Stadium, London – 31 May 2019 … Liverpool secured "
           "a dramatic victory…”", accent=CORAL)
quote_card(s, Inches(6.95), Inches(2.15), Inches(5.5), Inches(2.25),
           "Fine-tuned LED — 2001 FA Cup Final (true: Liverpool 2–1 Arsenal, Millennium Stadium)",
           "“Liverpool 1–0 Arsenal, Premier League, 20 May 2001 — Anfield, Anfield, "
           "Emirates Stadium, Anfield… Liverpool dominated the first half…”", accent=CORAL)
bullets(s, [
    [("Quantified over all 9 test matches:  ", {"bold": True, "color": CORAL}),
     ("fine-tuned BART (ROUGE rank 2) states the correct final score in ", {}),
     ("0 of 9 matches", {"bold": True, "color": CORAL}),
     (" and recalls only 30% of reference scorers. Fine-tuned LED (ROUGE rank 1): 2/9 scores, 13% scorers. "
      "BART scorelines are decorative — 1–1, 2–2, 1–0 regardless of the real 0–4, 3–0…", {})],
    [("ROUGE rewards the hallucination: ", {"bold": True}),
     ("the overlapping n-grams (team names, date formats, report style) are exactly what the "
      "models reproduce; wrong facts cost nothing.", {})],
], y=Inches(4.7), size=15)
footer(s, 8)

# --------------------------- 9 · Finding 4: the inversion
s = add_slide()
header(s, 9, "Finding 4", "Models optimise the reference, not the source",
       "BERTScore vs reference · NLI entailment vs transcript (DeBERTa-MNLI, TF-IDF-retrieved windows)")
figure_card(s, "figures/inversion.png", Inches(2.0), Inches(2.05), Inches(9.0))
bullets(s, [
    [("The rankings invert:  ", {"bold": True, "color": CORAL}),
     ("fine-tuned LED leads on ROUGE-L and BERTScore yet has only 2% NLI support; fine-tuned BART "
      "— zero NLI-supported sentences — is even more extreme. Prompted long-context LED, outranked "
      "on every reference metric, is the best-grounded condition (34%) because it copies the commentary.", {})],
    [("Fine-tuning on GPT-4 references optimises agreement with the references at the expense of "
      "grounding in the input. ", {"bold": True}),
     ("(NLI rates are conservative lower bounds — register gap, ASR noise — so compare across "
      "conditions, not absolutely.)", {"size": 13, "color": MUTED})],
], y=Inches(5.55), size=14)
footer(s, 9)

# ------------------------------ 10 · Data-quality incident + LED curve
s = add_slide()
header(s, 10, "Engineering", "The biggest improvement was not a model change",
       "Five silent LED training failures + one dataset bug")
bullets(s, [
    [("33% of training references were 0-byte files ", {"bold": True, "color": CORAL}),
     ("(and 44% of test references) — silently teaching models to emit empty strings and "
      "deflating every score ~1.8×.", {})],
    [("Predicted correction ratio matched the measured one almost exactly "
      "(e.g. fine-tuned BART 0.136 → 0.248 after the fix).", {})],
    [("LED also needed 5 bug fixes ", {"bold": True}),
     ("(defeated random-window sampler, collapse-to-EOS, decoder sentinel overflow, "
      "lost checkpoint, CLI truncation).", {})],
    [("After the fixes: a clean, monotone training curve → the earlier collapses were "
      "data and code defects, not modelling limits.", {"bold": True, "color": ACCENT})],
], x=Inches(0.85), y=Inches(2.15), w=Inches(6.2), size=15)
figure_card(s, "figures/led_training_curve.png", Inches(7.5), Inches(2.35), Inches(5.15))
tf = textbox(s, Inches(7.5), Inches(6.15), Inches(5.3), Inches(0.5))
set_run(tf.paragraphs[0].add_run(),
        "LED fine-tuning after the fixes: val ROUGE-L peaks at 0.172 (epoch 9).",
        11.5, MUTED, italic=True)
footer(s, 10)

# ---------------------------------------------- 11 · Conclusions
s = add_slide()
header(s, 11, "Conclusions", "What we learned, and where it breaks")
bullets(s, [
    [("1.  Fine-tuning a distribution-matched chunk pipeline is the best recipe at small N ", {"bold": True}),
     ("— +69% over the same pre-trained model, +55% over the best prompting baseline.", {})],
    [("2.  Context length is not a substitute for supervision ", {"bold": True}),
     ("— long-context wins only when no fine-tuning is possible.", {})],
    [("3.  Reference similarity and source grounding fully decouple ", {"bold": True}),
     ("— fine-tuned LED tops ROUGE-L and BERTScore yet achieves 2/9 correct scores and 2% NLI; "
      "BART goes further at 0/9 and 0% NLI. The transcript-copying prompted LED both outrank "
      "on ROUGE remains the best-grounded (34%). Optimising faithfulness, not just measuring it, is the next step.", {})],
    "",
    [("Limitations we state openly", {"bold": True, "color": CORAL})],
    [("Search-grounded GPT-4 references → facts anchored in web reports, but still ungrounded in the "
      "transcript, stylistically homogeneous, and not exactly reproducible (search results drift).",
      {"level": 1})],
    [("n = 9 test matches → bootstrap CIs resolve fine-tuning vs prompting (significant) but not "
      "LED vs BART (−0.0004 [−0.042, +0.035]).", {"level": 1})],
    [("N = 80 training pairs → both fine-tuned models overfit (val 0.31 → test 0.25 for BART).",
      {"level": 1})],
    [("NLI support rates are conservative lower bounds (ASR-noisy premises, register gap); "
      "entity precision is still unmeasured over all entity types.", {"level": 1})],
], size=16)
footer(s, 11)

os.makedirs("slides", exist_ok=True)
out = "slides/ANLP_slides_alt.pptx"
prs.save(out)
print("wrote", out)
